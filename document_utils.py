"""
Document processing utilities for extracting text from various file formats.
"""
import os
import re
import io

def read_txt(file_path):
    """Read text from a .txt file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        return f"Error reading TXT file: {str(e)}"

def read_docx(file_path):
    """Read text from a .docx file."""
    try:
        from docx import Document
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except ImportError:
        return "python-docx library not installed. Install with: pip install python-docx"
    except Exception as e:
        return f"Error reading DOCX file: {str(e)}"

def read_pptx(file_path):
    """Read text from a .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except ImportError:
        return "python-pptx library not installed. Install with: pip install python-pptx"
    except Exception as e:
        return f"Error reading PPTX file: {str(e)}"

def read_xlsx(file_path):
    """Read text from a .xlsx file without using the Excel agent."""
    try:
        import pandas as pd
        # Read all sheets
        xl = pd.ExcelFile(file_path)
        sheets = xl.sheet_names
        
        result = []
        result.append(f"File contains {len(sheets)} sheets: {', '.join(sheets)}")
        
        # Read each sheet into a dataframe and convert to text
        for sheet in sheets:
            df = pd.read_excel(file_path, sheet_name=sheet)
            result.append(f"\n\nSheet: {sheet}")
            result.append(f"Dimensions: {df.shape[0]} rows x {df.shape[1]} columns")
            result.append("Data Sample:")
            result.append(df.head(10).to_string())
        
        return '\n'.join(result)
    except ImportError:
        return "pandas library not installed. Install with: pip install pandas"
    except Exception as e:
        return f"Error reading XLSX file: {str(e)}"

def read_pdf(file_path):
    """Read text from a .pdf file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(file_path)
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return '\n'.join(text)
    except ImportError:
        return "PyPDF2 library not installed. Install with: pip install PyPDF2"
    except Exception as e:
        return f"Error reading PDF file: {str(e)}"

def read_csv(file_path):
    """Read text from a .csv file."""
    try:
        import pandas as pd
        df = pd.read_csv(file_path)
        result = []
        result.append(f"Dimensions: {df.shape[0]} rows x {df.shape[1]} columns")
        result.append("Data Sample:")
        result.append(df.head(10).to_string())
        return '\n'.join(result)
    except ImportError:
        return "pandas library not installed. Install with: pip install pandas"
    except Exception as e:
        return f"Error reading CSV file: {str(e)}"

def read_document(file_path):
    """
    Extract text from various document formats based on file extension.
    
    Args:
        file_path: Path to the document file
        
    Returns:
        Extracted text from the document
    """
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.txt':
        return read_txt(file_path)
    elif file_extension == '.docx':
        return read_docx(file_path)
    elif file_extension == '.pptx':
        return read_pptx(file_path)
    elif file_extension == '.xlsx' or file_extension == '.xls':
        return read_xlsx(file_path)
    elif file_extension == '.pdf':
        return read_pdf(file_path)
    elif file_extension == '.csv':
        return read_csv(file_path)
    else:
        return f"Unsupported file format: {file_extension}"

def read_document_bytes(file_bytes, file_name):
    """
    Extract text from document bytes based on file extension.
    
    Args:
        file_bytes: Bytes content of the file
        file_name: Name of the file (used to determine extension)
        
    Returns:
        Extracted text from the document
    """
    file_extension = os.path.splitext(file_name)[1].lower()
    
    # Create a temporary file
    temp_path = f"temp_{os.path.basename(file_name)}"
    try:
        with open(temp_path, "wb") as f:
            f.write(file_bytes)
        
        # Process the file
        result = read_document(temp_path)
        
        # Clean up
        if os.path.exists(temp_path):
            os.remove(temp_path)
            
        return result
    except Exception as e:
        if os.path.exists(temp_path):
            os.remove(temp_path)
        return f"Error processing file: {str(e)}"